# -*- coding: utf-8 -*-
import json
import textwrap
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT.parent / "weekly_report_20260715"
ASSET_DIR = OUT_DIR / "assets"
PPTX_PATH = OUT_DIR / "食物体积估计小程序后端链路阶段汇报_20260715_v2.pptx"
SCRIPT_PATH = OUT_DIR / "食物体积估计小程序后端链路阶段汇报_口头报告_20260715_v2.md"
WORK_DATASET = Path("D:/ .PROJECTS/CV-FOOD-DETECT/2026-4-3/VolETA-MetaFood/work".replace("D:/ ", "D:/"))

FONT = "Microsoft YaHei"
FONT_BOLD = "Microsoft YaHei UI"
BG = RGBColor(247, 249, 252)
NAVY = RGBColor(20, 37, 66)
BLUE = RGBColor(38, 100, 184)
CYAN = RGBColor(33, 145, 184)
GREEN = RGBColor(50, 142, 104)
GRAY = RGBColor(92, 104, 118)
LIGHT = RGBColor(226, 232, 240)
WHITE = RGBColor(255, 255, 255)
RED = RGBColor(185, 70, 70)


def ensure_dirs():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)


def load_json(path):
    return json.loads(Path(path).read_text(encoding="utf-8"))


def crop_image(src, dst, box):
    img = Image.open(src).convert("RGB")
    w, h = img.size
    left, top, right, bottom = box
    crop = img.crop((int(w * left), int(h * top), int(w * right), int(h * bottom)))
    crop.save(dst)
    return dst


def make_terminal_image(text, dst, width=1320, font_size=28, title=None):
    font_path = "C:/Windows/Fonts/consola.ttf"
    title_font_path = "C:/Windows/Fonts/msyhbd.ttc"
    font = ImageFont.truetype(font_path, font_size)
    title_font = ImageFont.truetype(title_font_path, 30)
    lines = []
    for raw in text.splitlines():
        wrapped = textwrap.wrap(raw, width=74, replace_whitespace=False) or [""]
        lines.extend(wrapped)
    line_h = font_size + 10
    pad = 34
    title_h = 46 if title else 0
    height = pad * 2 + title_h + line_h * len(lines)
    img = Image.new("RGB", (width, max(height, 260)), (22, 28, 38))
    draw = ImageDraw.Draw(img)
    y = pad
    if title:
        draw.text((pad, y), title, fill=(226, 232, 240), font=title_font)
        y += title_h
    for line in lines:
        color = (220, 230, 240)
        if "true" in line.lower() or "ok" in line.lower():
            color = (125, 211, 162)
        if "timeout" in line.lower():
            color = (245, 196, 93)
        draw.text((pad, y), line, fill=color, font=font)
        y += line_h
    img.save(dst)
    return dst


def make_code_image(text, dst, title):
    return make_terminal_image(text, dst, width=1380, font_size=24, title=title)


def dataset_summary():
    dirs = [p for p in WORK_DATASET.iterdir() if p.is_dir()]
    rows = []
    total_jpg = 0
    total_png = 0
    first_images = []
    for folder in sorted(dirs, key=lambda p: "".join(ch for ch in p.name if ch.isdigit()).zfill(4)):
        files = list(folder.rglob("*"))
        jpgs = [p for p in files if p.suffix.lower() == ".jpg"]
        pngs = [p for p in files if p.suffix.lower() == ".png"]
        if jpgs or pngs:
            first = sorted(jpgs or pngs)[0]
            first_images.append((folder.name, first))
        total_jpg += len(jpgs)
        total_png += len(pngs)
        if jpgs or pngs:
            rows.append((folder.name, len(jpgs) + len(pngs)))
    return {
        "sample_dirs": len([r for r in rows if r[1] > 0]),
        "total_jpg": total_jpg,
        "total_png": total_png,
        "rows": rows,
        "first_images": first_images,
    }


def make_dataset_montage(dst):
    summary = dataset_summary()
    items = summary["first_images"][:12]
    thumb_w, thumb_h = 260, 185
    gap = 18
    label_h = 34
    cols = 4
    rows = 3
    width = cols * thumb_w + (cols + 1) * gap
    height = rows * (thumb_h + label_h) + (rows + 1) * gap
    img = Image.new("RGB", (width, height), (247, 249, 252))
    draw = ImageDraw.Draw(img)
    font = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 22)
    for idx, (name, path) in enumerate(items):
        r, c = divmod(idx, cols)
        x = gap + c * (thumb_w + gap)
        y = gap + r * (thumb_h + label_h + gap)
        src = Image.open(path).convert("RGB")
        src.thumbnail((thumb_w, thumb_h))
        canvas = Image.new("RGB", (thumb_w, thumb_h), (226, 232, 240))
        ox = (thumb_w - src.width) // 2
        oy = (thumb_h - src.height) // 2
        canvas.paste(src, (ox, oy))
        img.paste(canvas, (x, y))
        draw.text((x, y + thumb_h + 5), f"work/{name}", fill=(20, 37, 66), font=font)
    img.save(dst)
    return dst


def set_font(run, size=18, color=NAVY, bold=False):
    run.font.name = FONT_BOLD if bold else FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_box(slide, x, y, w, h, fill=WHITE, line=LIGHT, radius=True):
    shape = slide.shapes.add_shape(5 if radius else 1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def add_text(slide, text, x, y, w, h, size=18, color=NAVY, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size, color, bold)
    return box


def add_bullets(slide, items, x, y, w, h, size=16, color=NAVY):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(7)
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return box


def add_header(slide, title, subtitle=None):
    add_text(slide, title, 0.55, 0.32, 11.0, 0.45, size=25, color=NAVY, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.58, 0.83, 11.4, 0.35, size=11.5, color=GRAY)
    line = slide.shapes.add_shape(1, Inches(0.55), Inches(1.18), Inches(12.2), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT
    line.line.color.rgb = LIGHT


def add_footer(slide, page):
    add_text(slide, "食物体积估计小程序后端链路阶段汇报", 0.58, 7.15, 6.4, 0.2, size=8.5, color=GRAY)
    add_text(slide, f"{page:02d}", 12.35, 7.15, 0.35, 0.2, size=8.5, color=GRAY, align=PP_ALIGN.RIGHT)


def add_image(slide, path, x, y, w, h):
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def flow(slide, labels, x, y, box_w, box_h, gap):
    for i, label in enumerate(labels):
        bx = x + i * (box_w + gap)
        add_box(slide, bx, y, box_w, box_h, fill=WHITE, line=RGBColor(190, 205, 225))
        add_text(slide, label, bx + 0.08, y + 0.18, box_w - 0.16, box_h - 0.25, size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        if i < len(labels) - 1:
            add_text(slide, "→", bx + box_w + 0.04, y + 0.22, gap - 0.08, box_h - 0.2, size=22, color=BLUE, bold=True, align=PP_ALIGN.CENTER)


def table(slide, headers, rows, x, y, w, h):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = shape.table
    for col, header in enumerate(headers):
        cell = tbl.cell(0, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.text = header
        p.font.name = FONT
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(248, 250, 252) if r % 2 else WHITE
            p = cell.text_frame.paragraphs[0]
            p.text = value
            p.font.name = FONT
            p.font.size = Pt(10.5)
            p.font.color.rgb = NAVY
    return shape


def build_assets():
    startup_src = Path("C:/Users/AW/AppData/Local/Temp/codex-clipboard-816180b0-dcec-43c9-b412-956749c3845a.png")
    install_src = Path("C:/Users/AW/AppData/Local/Temp/codex-clipboard-114370fe-8525-40c2-bfc9-74e6d9631163.png")
    if startup_src.exists():
        crop_image(startup_src, ASSET_DIR / "terminal_uvicorn_started.png", (0.0, 0.48, 1.0, 0.98))
    if install_src.exists():
        crop_image(install_src, ASSET_DIR / "terminal_install_warning.png", (0.0, 0.0, 1.0, 0.76))

    health = """Invoke-RestMethod http://127.0.0.1:8000/health
{
  "status": "ok",
  "checkpoint_exists": true,
  "checkpoint": "D:\\\\.PROJECTS\\\\CV-FOOD-DETECT\\\\2026-7-03\\\\miniapp_food_volume_bundle\\\\models\\\\depth_pro.pt",
  "timeout_seconds": 10.0
}"""
    make_terminal_image(health, ASSET_DIR / "terminal_health.png", title="后端健康检查返回")

    endpoint = """@app.get("/health")
def health() -> Dict[str, Any]:
    checkpoint = _checkpoint_path()
    return {
        "status": "ok",
        "checkpoint_exists": checkpoint.exists(),
        "checkpoint": str(checkpoint),
        "timeout_seconds": _timeout_seconds(),
    }

@app.post("/api/v1/food-volume")
async def estimate_food_volume(
    image: UploadFile = File(...),
    food: Optional[str] = Form(None),
    phone_model: Optional[str] = Form(None),
    prompt: Optional[str] = Form(None),
    mask: Optional[UploadFile] = File(None),
    intrinsics: Optional[UploadFile] = File(None),
) -> Dict[str, Any]:"""
    make_code_image(endpoint, ASSET_DIR / "code_backend_api.png", "backend/app.py：云端接口封装")

    timing = """"timings_seconds": {
  "auto_mask_seconds": 0.0,
  "depthpro_seconds": <measured>,
  "volume_seconds": <measured>,
  "total_seconds": <measured>
}"""
    make_code_image(timing, ASSET_DIR / "code_timings.png", "pipeline_summary.json：分段耗时字段")

    tree = """miniapp_food_volume_bundle/
├─ backend/
│  ├─ app.py
│  └─ README.md
├─ scripts/
│  ├─ run_auto_food_volume_demo.py
│  ├─ run_depthpro_demo.py
│  └─ estimate_metric_volume_calorie.py
├─ docs/
│  └─ 前后端边界与云端部署.md
├─ reports/
│  ├─ week_20260715_plan.md
│  └─ timing_20260715.md
└─ models/
   └─ depth_pro.pt"""
    make_terminal_image(tree, ASSET_DIR / "project_tree.png", title="本周整理后的交付结构")
    make_dataset_montage(ASSET_DIR / "work_dataset_montage.png")


def build_ppt():
    demo = load_json(ROOT / "runs/depthpro_smoke_test/demo_summary.json")
    scaffold = load_json(ROOT / "runs/depthpro_smoke_test/02_results/metric_scaffold/summary.json")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def blank():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BG
        return slide

    # 1
    slide = blank()
    add_box(slide, 0.0, 0.0, 13.333, 7.5, fill=NAVY, line=NAVY, radius=False)
    add_text(slide, "食物体积估计小程序后端链路阶段汇报", 0.75, 1.2, 11.8, 0.7, size=33, color=WHITE, bold=True)
    add_text(slide, "云端推理约束、前后端职责拆分、链路优化与模型调研", 0.78, 2.05, 11.5, 0.45, size=18, color=RGBColor(210, 222, 240))
    add_box(slide, 0.78, 3.05, 4.05, 1.35, fill=RGBColor(32, 55, 92), line=RGBColor(64, 100, 150))
    add_text(slide, "本周实际调整", 1.05, 3.28, 3.5, 0.25, size=15, color=RGBColor(180, 215, 255), bold=True)
    add_text(slide, "模型推理移到云端；本周完成后端入口、健康检查、模型路径检查和耗时字段记录。", 1.05, 3.68, 3.45, 0.55, size=14, color=WHITE)
    add_text(slide, "2026.07.15", 0.82, 6.55, 3.0, 0.35, size=15, color=RGBColor(210, 222, 240))

    # 2
    slide = blank()
    add_header(slide, "本周工作记录")
    cards = [
        ("部署约束调整", "模型部分放在云端；本周先验证服务入口、模型路径、超时参数和接口返回。"),
        ("链路重新确认", "保留单图 depth + support-plane volume 路线，减少前端输入，后端负责完整推理。"),
        ("后续模型对比", "列出可替换模型，但不直接替换主线；后续使用 work 数据集做同图测试。"),
    ]
    for i, (title, body) in enumerate(cards):
        x = 0.8 + i * 4.1
        add_box(slide, x, 1.65, 3.65, 3.55, fill=WHITE, line=RGBColor(205, 214, 226))
        add_text(slide, f"0{i+1}", x + 0.22, 1.9, 0.7, 0.35, size=18, color=BLUE, bold=True)
        add_text(slide, title, x + 0.9, 1.92, 2.45, 0.35, size=18, color=NAVY, bold=True)
        add_text(slide, body, x + 0.35, 2.65, 2.95, 1.5, size=15, color=GRAY)
    add_footer(slide, 2)

    # 3
    slide = blank()
    add_header(slide, "总体架构调整：端侧轻量化，云端集中推理")
    add_bullets(
        slide,
        [
            "小程序端保留拍摄、上传、可选食物类别与展示功能。",
            "云端后端承担自动 mask、相机参数推断、metric depth、3D scaffold、体积与热量估计。",
            "服务侧新增统一 HTTP 接口，避免小程序直接调用算法脚本。",
            "后续性能评估从端侧模型体量转为云端热启动耗时与服务稳定性。"
        ],
        0.75, 1.55, 5.2, 3.8, size=16
    )
    add_box(slide, 6.35, 1.45, 5.9, 4.55, fill=WHITE, line=RGBColor(205, 214, 226))
    add_text(slide, "前端", 6.7, 1.8, 1.2, 0.35, size=17, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "拍摄/上传图片\n可选食物类别\n可选手机型号\n结果展示", 6.55, 2.35, 1.5, 1.6, size=13, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, "→", 8.15, 2.65, 0.45, 0.35, size=24, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "云端后端", 8.75, 1.8, 1.4, 0.35, size=17, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "FastAPI 接口\nDepthPro 推理\n3D scaffold\n体积与 kcal", 8.6, 2.35, 1.7, 1.6, size=13, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, "→", 10.45, 2.65, 0.45, 0.35, size=24, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "输出", 11.05, 1.8, 1.0, 0.35, size=17, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "volume\nmass\nkcal\n可视化结果", 10.9, 2.35, 1.3, 1.6, size=13, color=NAVY, align=PP_ALIGN.CENTER)
    add_footer(slide, 3)

    # 4
    slide = blank()
    add_header(slide, "当前主链路：RGB 图像到体积与热量")
    flow(slide, ["RGB 图像", "自动 mask", "相机参数", "DepthPro", "3D scaffold", "支撑平面", "体积/kcal"], 0.45, 1.55, 1.45, 0.9, 0.18)
    add_box(slide, 0.75, 3.05, 5.6, 2.65, fill=WHITE, line=RGBColor(205, 214, 226))
    add_text(slide, "核心计算逻辑", 1.05, 3.32, 4.4, 0.3, size=17, color=NAVY, bold=True)
    add_bullets(slide, [
        "通过 metric depth 将图像像素反投影为三维点集。",
        "基于 mask 提取食物区域的 3D scaffold。",
        "利用食物周边区域拟合支撑平面，进行高度积分得到体积。",
        "结合食物密度与 kcal/100g 先验输出热量估计。"
    ], 1.05, 3.8, 4.8, 1.55, size=13.5)
    add_box(slide, 6.8, 3.05, 5.35, 2.65, fill=WHITE, line=RGBColor(205, 214, 226))
    add_text(slide, "本周代码变更", 7.1, 3.32, 4.4, 0.3, size=17, color=NAVY, bold=True)
    add_bullets(slide, [
        "后端服务接口：`/api/v1/food-volume`。",
        "健康检查接口：`/health`。",
        "推理超时控制：`FOOD_VOLUME_TIMEOUT_SECONDS`。",
        "分段耗时记录：mask、DepthPro、volume、total。"
    ], 7.1, 3.8, 4.7, 1.55, size=13.5)
    add_footer(slide, 4)

    # 5
    slide = blank()
    add_header(slide, "与前期 VoIETA / VoITEX 路线的差异")
    rows = [
        ["重建目标", "NeRF / mesh 完整重建", "单图 metric depth + 3D scaffold"],
        ["推理链路", "多图/视频、训练或重建过程较重", "单图输入，链路更短，适合后端接口化"],
        ["尺度来源", "重建坐标与后续定尺耦合", "DepthPro 焦距与 metric depth 支持无标定物尺度估计"],
        ["工程定位", "偏实验复现与 3D 重建验证", "偏小程序后端服务与体积计算闭环"],
        ["当前风险", "耗时长、部署复杂", "深度模型与分割质量决定体积稳定性"],
    ]
    table(slide, ["维度", "前期方案", "当前方案"], rows, 0.65, 1.45, 12.0, 4.85)
    add_footer(slide, 5)

    # 6
    slide = blank()
    add_header(slide, "代码规范化与服务化改动")
    add_image(slide, ASSET_DIR / "project_tree.png", 0.65, 1.35, 5.75, 4.85)
    add_bullets(slide, [
        "新增 `backend/app.py`，将命令行推理封装为 FastAPI 服务。",
        "新增 `backend/README.md`，固定请求字段、返回字段和启动方式。",
        "新增 `docs/前后端边界与云端部署.md`，明确小程序端与云端后端边界。",
        "新增 `reports/week_20260715_plan.md`，记录本周技术目标和模型调研方向。",
        "更新 `requirements.txt`，加入 FastAPI、Uvicorn 和 multipart 上传依赖。"
    ], 6.75, 1.65, 5.5, 3.7, size=15)
    add_footer(slide, 6)

    # 7
    slide = blank()
    add_header(slide, "后端接口设计")
    add_image(slide, ASSET_DIR / "code_backend_api.png", 0.65, 1.28, 6.2, 4.95)
    add_box(slide, 7.15, 1.45, 5.25, 4.55, fill=WHITE, line=RGBColor(205, 214, 226))
    add_text(slide, "接口输入", 7.45, 1.78, 2.0, 0.3, size=17, color=NAVY, bold=True)
    add_bullets(slide, [
        "`image`：必填，用户拍摄图片。",
        "`food`：可选，食物类别。",
        "`phone_model`：可选，相机参数匹配。",
        "`prompt`：可选，自动分割文本提示。",
        "`mask / intrinsics`：调试阶段可选输入。"
    ], 7.45, 2.25, 4.55, 1.55, size=13.5)
    add_text(slide, "接口输出", 7.45, 4.0, 2.0, 0.3, size=17, color=NAVY, bold=True)
    add_bullets(slide, [
        "`request_id`、`elapsed_seconds`、`output_dir`。",
        "`pipeline_summary`：链路摘要与分段耗时。",
        "`volume_summary`：体积、质量和热量估计。"
    ], 7.45, 4.45, 4.55, 1.0, size=13.5)
    add_footer(slide, 7)

    # 8
    slide = blank()
    add_header(slide, "10 秒推理目标的量化记录方式")
    add_image(slide, ASSET_DIR / "code_timings.png", 0.75, 1.4, 5.9, 2.8)
    add_box(slide, 7.0, 1.4, 5.2, 3.55, fill=WHITE, line=RGBColor(205, 214, 226))
    add_text(slide, "计时口径", 7.3, 1.75, 1.8, 0.3, size=17, color=NAVY, bold=True)
    add_bullets(slide, [
        "`auto_mask_seconds`：自动检测与分割耗时。",
        "`depthpro_seconds`：DepthPro 深度估计与 scaffold 构建耗时。",
        "`volume_seconds`：支撑平面体积与热量计算耗时。",
        "`total_seconds`：完整请求链路总耗时。"
    ], 7.3, 2.25, 4.5, 1.7, size=14)
    add_text(slide, "说明：正式性能评估需在云端 GPU 环境中进行热启动测试；本地 CPU 环境不作为最终性能结论。", 0.85, 5.2, 11.5, 0.6, size=13.5, color=GRAY)
    add_footer(slide, 8)

    # 9
    slide = blank()
    add_header(slide, "后端服务运行验证")
    add_image(slide, ASSET_DIR / "terminal_uvicorn_started.png", 0.65, 1.35, 6.2, 3.15)
    add_image(slide, ASSET_DIR / "terminal_health.png", 7.05, 1.35, 5.6, 2.55)
    add_box(slide, 7.05, 4.25, 5.6, 1.25, fill=WHITE, line=RGBColor(205, 214, 226))
    add_text(slide, "验证结果", 7.35, 4.5, 1.4, 0.25, size=16, color=NAVY, bold=True)
    add_bullets(slide, [
        "Uvicorn 服务已在 8000 端口启动。",
        "`/health` 返回 `status=ok`。",
        "`checkpoint_exists=true`，本地 DepthPro 权重路径可访问。"
    ], 7.35, 4.88, 4.8, 0.72, size=12.2)
    add_footer(slide, 9)

    # 10
    slide = blank()
    add_header(slide, "DepthPro 与 3D scaffold 验证输出")
    img_paths = [
        ROOT / "runs/depthpro_smoke_test/01_process/input_with_mask_bbox.jpg",
        ROOT / "runs/depthpro_smoke_test/01_process/depth_color_masked.png",
        ROOT / "runs/depthpro_smoke_test/02_results/metric_scaffold/overlay_projection.png",
        ROOT / "examples/output_support_plane/support_plane_overlay.png",
    ]
    labels = ["输入图像与食物区域", "DepthPro 深度可视化", "3D scaffold 投影", "支撑平面区域"]
    for i, (path, label) in enumerate(zip(img_paths, labels)):
        x = 0.65 + (i % 2) * 6.25
        y = 1.35 + (i // 2) * 2.65
        add_box(slide, x, y, 5.65, 2.25, fill=WHITE, line=RGBColor(205, 214, 226))
        add_image(slide, path, x + 0.08, y + 0.08, 2.55, 2.05)
        add_text(slide, label, x + 2.82, y + 0.35, 2.45, 0.35, size=15, color=NAVY, bold=True)
        add_text(slide, "用于验证图像输入、mask、metric depth 与 scaffold 生成链路。", x + 2.82, y + 0.86, 2.45, 0.8, size=11.5, color=GRAY)
    add_footer(slide, 10)

    # 11
    slide = blank()
    add_header(slide, "本地 smoke test 摘要")
    bbox = scaffold["bbox_camera"]["extent"]
    rows = [
        ["模型分支", "Apple DepthPro"],
        ["焦距预测", f"{demo['predicted_focallength_px']:.2f} px"],
        ["食物区域深度范围", f"{demo['masked_depth_stats_m']['min']:.3f} m - {demo['masked_depth_stats_m']['max']:.3f} m"],
        ["scaffold 点数", f"{scaffold['counts']['food_scaffold_points']}"],
        ["bbox extent", f"{bbox[0]:.3f} m × {bbox[1]:.3f} m × {bbox[2]:.3f} m"],
    ]
    table(slide, ["指标", "结果"], rows, 0.85, 1.55, 5.75, 3.25)
    add_box(slide, 7.0, 1.55, 5.4, 3.25, fill=WHITE, line=RGBColor(205, 214, 226))
    add_text(slide, "解释边界", 7.3, 1.9, 1.8, 0.3, size=17, color=NAVY, bold=True)
    add_bullets(slide, [
        "该测试用于证明 DepthPro、mask、scaffold 输出链路可执行。",
        "当前样例不作为最终体积精度结论。",
        "后续基于已有 work 数据集统计误差与失败类型。",
        "复杂中餐场景仍需重点验证分割与支撑平面鲁棒性。"
    ], 7.3, 2.35, 4.55, 1.8, size=14)
    add_footer(slide, 11)

    # 12
    dataset = dataset_summary()
    slide = blank()
    add_header(slide, "已有数据集情况")
    add_image(slide, ASSET_DIR / "work_dataset_montage.png", 0.55, 1.35, 6.45, 4.45)
    add_box(slide, 7.35, 1.45, 4.9, 3.95, fill=WHITE, line=RGBColor(205, 214, 226))
    add_text(slide, "本地数据目录", 7.65, 1.78, 2.2, 0.3, size=17, color=NAVY, bold=True)
    add_text(slide, str(WORK_DATASET), 7.65, 2.18, 4.25, 0.55, size=11.2, color=GRAY)
    add_bullets(slide, [
        f"有效编号样例目录：{dataset['sample_dirs']} 个。",
        f"图像文件：{dataset['total_jpg']} 张 jpg，{dataset['total_png']} 张 png。",
        "目录中包含多帧图像、单图样例和前期重建中间结果。",
        "下一步不是重新建立数据集，而是从该目录筛选固定测试子集。"
    ], 7.65, 2.95, 4.2, 1.55, size=13.5)
    add_text(slide, "后续测试口径：固定若干样例，分别记录 mask、depth、support-plane、volume 的耗时与失败原因。", 7.65, 4.72, 4.15, 0.45, size=12.5, color=GRAY)
    add_footer(slide, 12)

    # 13
    slide = blank()
    add_header(slide, "可替代模型与类似项目调研")
    rows = [
        ["DepthPro", "metric depth + 焦距估计", "当前主线；适合云端验证"],
        ["UniDepth / UniDepthV2", "通用 metric depth", "作为无内参尺度恢复对照"],
        ["Metric3D / Metric3D V2", "通用单目 metric depth", "作为深度模型横向对比"],
        ["Depth Anything V2 Metric", "模型生态完整、体量选择更多", "需验证近距离食物尺度偏差"],
        ["Grounded-SAM-2 / SAM 2", "文本提示或交互式分割", "提升自动 mask 稳定性"],
        ["YOLOE / 开放词表分割", "检测与分割一体化", "减少 GroundingDINO 依赖"],
    ]
    table(slide, ["候选方向", "作用", "本项目定位"], rows, 0.55, 1.35, 12.25, 4.85)
    add_footer(slide, 13)

    # 14
    slide = blank()
    add_header(slide, "下一阶段计划")
    items = [
        ("性能基准", "在云端 GPU 环境中记录冷启动与热启动耗时，目标为完整推理约 10 秒。"),
        ("模型对比", "选择 UniDepth、Metric3D 或 Depth Anything V2 Metric 之一与 DepthPro 进行同图对照。"),
        ("数据集使用", "从 `work` 目录筛选固定测试子集，保留输入图、输出图、耗时和失败原因。"),
        ("鲁棒性优化", "优先优化盘子/桌面支撑平面识别，并扩充中式菜肴密度与热量先验。"),
    ]
    for i, (title, body) in enumerate(items):
        x = 0.75 + (i % 2) * 6.1
        y = 1.55 + (i // 2) * 2.25
        add_box(slide, x, y, 5.4, 1.6, fill=WHITE, line=RGBColor(205, 214, 226))
        add_text(slide, title, x + 0.32, y + 0.25, 1.8, 0.3, size=17, color=BLUE, bold=True)
        add_text(slide, body, x + 0.32, y + 0.75, 4.75, 0.52, size=13.5, color=GRAY)
    add_footer(slide, 14)

    prs.save(PPTX_PATH)


def build_script():
    text = """# 食物体积估计小程序后端链路阶段汇报口头报告

老师好，本周工作主要围绕小程序后端对接和技术链路重新整理展开。根据本周新的要求，模型部分可以放在云端，因此当前工作重点不再是压缩到小程序端侧运行，而是保证云端能够稳定完成图片三维信息恢复、体积计算和热量估计，并尽量将单次图片推理控制在十秒左右。

第一部分是工程约束的变化。此前我们一直关注模型体量是否能直接放入小程序环境，但现在模型可以部署在服务器端，因此前端只需要完成拍摄、上传、可选食物类别输入和结果展示。深度估计、自动分割、三维点支架构建、体积计算和热量估计都交给后端完成。这样可以降低小程序端的部署压力，也便于后续替换模型和记录运行日志。

第二部分是本周对代码结构进行的整理。当前已经新增了 `backend/app.py`，将原来的命令行推理流程封装成 FastAPI 后端服务。后端提供两个主要接口：一个是 `/health`，用于检查服务状态和 DepthPro 权重路径；另一个是 `/api/v1/food-volume`，用于接收小程序上传的图片，并返回体积、质量、热量和中间结果摘要。当前健康检查已经返回 `status=ok`，并且 `checkpoint_exists=true`，说明后端服务和模型路径已经能够正常识别。

第三部分是前后端边界的明确。小程序前端只保留必要交互：上传 RGB 图像、可选上传食物类别、可选上传手机型号，并展示后端返回结果。云端后端负责自动生成 mask、推断相机参数、调用 DepthPro 得到 metric depth、构建 3D scaffold、基于支撑平面计算体积，并结合密度和热量先验输出 kcal。这样可以避免让前端直接接触算法脚本，也方便后续师兄进行小程序接口集成。

第四部分是当前技术链路。整体流程是：输入 RGB 图像后，先通过自动分割得到食物 mask；随后根据手机型号、EXIF 或 DepthPro 的焦距预测推断相机内参；然后使用 DepthPro 生成绝对尺度深度图；再将深度图和 mask 反投影为三维点集，形成 3D scaffold；最后利用食物周围区域拟合支撑平面，通过高度积分计算体积，并结合食物密度和 kcal/100g 先验得到热量估计。

第五部分是与前期 VoIETA 和 VoITEX 方案的区别。前期方案更偏向 NeRF 或 mesh 的完整三维重建，链路较重，更适合验证三维重建能力。当前方案不再把完整 mesh 重建作为主路径，而是使用单图 metric depth 和支撑平面直接完成体积计算。这样链路更短，推理更容易服务化，也更符合小程序后端接入的需要。相应地，当前方案的主要风险转为深度模型质量、自动分割质量和支撑平面拟合稳定性。

第六部分是运行验证。本周已经完成依赖安装和后端服务启动。Uvicorn 服务能够在 8000 端口启动，健康检查接口能够返回模型路径存在。PPT 中放入了服务启动截图、健康检查结果，以及 DepthPro smoke test 的输出图，包括输入图像与 mask、深度可视化、scaffold 投影和支撑平面区域。这些结果说明当前链路已经具备后端服务化验证基础。

第七部分是耗时统计方式。本周在 `run_auto_food_volume_demo.py` 中新增了分段耗时字段，包括 `auto_mask_seconds`、`depthpro_seconds`、`volume_seconds` 和 `total_seconds`。后续只要在云端环境中运行，即可直接得到每一阶段的耗时，用于评估是否达到十秒左右的目标。本机当前 Python 环境虽然可以启动服务，但不是完整 GPU 推理环境，因此本地 CPU 测试结果不作为最终性能结论。

第八部分是模型和类似项目调研。深度模型方面，当前主线仍然是 Apple DepthPro，因为它能够预测焦距，比较适合无标定物和无相机内参输入的场景。备选方向包括 UniDepth、Metric3D 和 Depth Anything V2 Metric，后续可以选择其中一个和 DepthPro 做同图对照。分割模型方面，当前可以继续沿用 GroundingDINO 与 SAM 的思路，也可以调研 Grounded-SAM-2、SAM 2 或开放词表检测分割模型，目标是进一步减少用户输入，提高复杂中式菜肴的 mask 稳定性。

第九部分是数据集情况。当前本地 `work` 目录已经有可用数据，不需要重新建立样例集。我检查到该目录下有二十个有效编号样例目录，包含三千八百多张 jpg 图像和接近两百张 png 图像，其中既有多帧图像，也有单图样例和前期重建过程中的中间结果。下一步应当从这里筛选固定测试子集，用同一批数据记录分割、深度估计、支撑平面计算和体积输出的耗时与失败原因。

下一阶段计划主要有四项。第一，在云端 GPU 环境中建立正式性能基准，统计冷启动和热启动耗时。第二，选择一个新的 metric depth 模型与 DepthPro 进行横向对比。第三，从已有 `work` 数据集中筛选固定测试子集，保留输入图、输出图、耗时和失败原因。第四，优先优化桌面和盘子支撑平面识别，并扩充中式菜肴的密度和热量先验库。
"""
    SCRIPT_PATH.write_text(text, encoding="utf-8")


def main():
    ensure_dirs()
    build_assets()
    build_ppt()
    build_script()
    print(PPTX_PATH)
    print(SCRIPT_PATH)


if __name__ == "__main__":
    main()
